"""Extract bounded text from common document formats in an isolated process."""

import csv
import email
from email import policy
import subprocess
import sys
import zipfile
from pathlib import Path


ZIP_BASED_EXTENSIONS = {
    ".docx", ".docm", ".dotx", ".dotm",
    ".xlsx", ".xlsm", ".xltx", ".xltm", ".xlsb",
    ".pptx", ".pptm", ".ppsx", ".potx", ".potm",
    ".odt", ".ods", ".odp", ".odg", ".epub",
}
MAX_ARCHIVE_MEMBERS = 10_000
MAX_UNCOMPRESSED_BYTES = 512 * 1024 * 1024


def apply_process_limits():
    """Bound hostile/accidental document expansion on supported Unix hosts."""
    try:
        import resource

        resource.setrlimit(resource.RLIMIT_AS, (384 * 1024 * 1024, 384 * 1024 * 1024))
        resource.setrlimit(resource.RLIMIT_CPU, (15, 15))
    except (ImportError, OSError, ValueError):
        # Windows development and some containers do not expose these limits.
        pass


def validate_zip_container(path):
    """Reject damaged files and zip bombs before an Office/ODF/EPUB parser runs."""
    with zipfile.ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > MAX_ARCHIVE_MEMBERS:
            raise ValueError("Document contains too many compressed members")
        total_size = sum(member.file_size for member in members)
        if total_size > MAX_UNCOMPRESSED_BYTES:
            raise ValueError("Document expands beyond the safe extraction limit")


class Collector:
    def __init__(self, limit):
        self.limit = limit
        self.parts = []
        self.size = 0

    @property
    def full(self):
        return self.size >= self.limit

    def add(self, value):
        if self.full or value is None:
            return
        text = str(value).strip()
        if not text:
            return
        remaining = self.limit - self.size
        fragment = text[:remaining]
        self.parts.append(fragment)
        self.size += len(fragment)

    def output(self):
        return "\n".join(self.parts)[:self.limit]


def extract_docx(path, collector):
    from docx import Document

    document = Document(path)
    for paragraph in document.paragraphs:
        collector.add(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            collector.add(" | ".join(cell.text for cell in row.cells))


def extract_xlsx(path, collector):
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            collector.add(sheet.title)
            for row in sheet.iter_rows(values_only=True):
                collector.add(" | ".join(str(value) for value in row if value is not None))
                if collector.full:
                    return
    finally:
        workbook.close()


def extract_xls(path, collector):
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets():
            collector.add(sheet.name)
            for row_index in range(sheet.nrows):
                collector.add(" | ".join(str(value) for value in sheet.row_values(row_index)))
                if collector.full:
                    return
    finally:
        workbook.release_resources()


def extract_xlsb(path, collector):
    from pyxlsb import open_workbook

    with open_workbook(path) as workbook:
        for sheet_name in workbook.sheets:
            collector.add(sheet_name)
            with workbook.get_sheet(sheet_name) as sheet:
                for row in sheet.rows():
                    collector.add(" | ".join(str(cell.v) for cell in row if cell.v is not None))
                    if collector.full:
                        return


def extract_pptx(path, collector):
    from pptx import Presentation

    presentation = Presentation(path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                collector.add(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    collector.add(" | ".join(cell.text for cell in row.cells))
            if collector.full:
                return


def extract_odf(path, collector):
    from odf import teletype, text
    from odf.opendocument import load

    document = load(path)
    for element_type in (text.H, text.P):
        for element in document.getElementsByType(element_type):
            collector.add(teletype.extractText(element))
            if collector.full:
                return


def extract_epub(path, collector):
    from bs4 import BeautifulSoup
    from ebooklib import ITEM_DOCUMENT, epub

    book = epub.read_epub(path, options={"ignore_ncx": True})
    for title, _attributes in book.get_metadata("DC", "title"):
        collector.add(title)
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        collector.add(BeautifulSoup(item.get_content(), "html.parser").get_text(" ", strip=True))
        if collector.full:
            return


def extract_eml(path, collector):
    with open(path, "rb") as source:
        message = email.message_from_binary_file(source, policy=policy.default)
    for header in ("subject", "from", "to", "date"):
        collector.add(message.get(header))
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                collector.add(part.get_content())
    else:
        collector.add(message.get_content())


def extract_legacy(path, collector, command):
    result = subprocess.run(
        [command, str(path)],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=8,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError((result.stderr or f"{command} failed")[:300])
    collector.add(result.stdout)


def extract_rtf(path, collector):
    from striprtf.striprtf import rtf_to_text

    collector.add(rtf_to_text(path.read_text(encoding="utf-8", errors="replace")))


def extract_delimited(path, collector, delimiter):
    with path.open("r", encoding="utf-8", errors="replace", newline="") as source:
        for row in csv.reader(source, delimiter=delimiter):
            collector.add(" | ".join(row))
            if collector.full:
                return


def main():
    apply_process_limits()
    path = Path(sys.argv[1])
    collector = Collector(int(sys.argv[2]))
    suffix = path.suffix.lower()
    if suffix in ZIP_BASED_EXTENSIONS:
        validate_zip_container(path)
    if suffix in {".docx", ".docm", ".dotx", ".dotm"}:
        extract_docx(path, collector)
    elif suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        extract_xlsx(path, collector)
    elif suffix == ".xls":
        extract_xls(path, collector)
    elif suffix == ".xlsb":
        extract_xlsb(path, collector)
    elif suffix in {".pptx", ".pptm", ".ppsx", ".potx", ".potm"}:
        extract_pptx(path, collector)
    elif suffix in {".odt", ".ods", ".odp", ".odg"}:
        extract_odf(path, collector)
    elif suffix == ".rtf":
        extract_rtf(path, collector)
    elif suffix == ".epub":
        extract_epub(path, collector)
    elif suffix == ".eml":
        extract_eml(path, collector)
    elif suffix == ".doc":
        extract_legacy(path, collector, "antiword")
    elif suffix == ".ppt":
        extract_legacy(path, collector, "catppt")
    else:
        raise ValueError(f"Unsupported document extension: {suffix}")
    output = collector.output().strip()
    if not output:
        return 2
    sys.stdout.write(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
